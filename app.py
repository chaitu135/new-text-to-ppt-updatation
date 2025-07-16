import streamlit as st
import base64
import os
from pptx import Presentation
from pptx.util import Pt
from transformers import pipeline

# Format settings
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def get_generator():
    # Use Streamlit resource cache for efficiency
    @st.cache_resource
    def load_generator():
        return pipeline("text-generation", model="gpt2")
    return load_generator()

generator = get_generator()

def generate_slide_titles(topic: str, openai_api_key: str = None) -> list:
    """
    Generate informative PowerPoint slide titles for a given topic using OpenAI API or fallback
    """
    def fallback_titles():
        prompt = f"Generate 5 professional PowerPoint slide titles about: {topic}\n1."
        result = generator(
            prompt,
            max_new_tokens=80,
            num_return_sequences=1,
            truncation=True,
            pad_token_id=50256
        )[0]['generated_text']
        lines = result.split("\n")
        titles = [line.strip("-•1234567890. ") for line in lines if line.strip()]
        return [t for t in titles if t][:5]
    if openai_api_key:
        try:
            import openai
            client = openai.OpenAI(api_key=openai_api_key)
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a professional PowerPoint presentation assistant. Generate 5 concise slide titles."},
                    {"role": "user", "content": f"Generate 5 professional PowerPoint slide titles about: {topic}"}
                ],
                temperature=0.7,
                max_tokens=200
            )
            content = response.choices[0].message.content
            titles = [line.strip("-•1234567890. ") for line in content.split("\n") if line.strip()]
            titles = [t for t in titles if t][:5]
            if not titles:
                st.warning("OpenAI did not return any titles, using fallback.")
                return fallback_titles()
            return titles
        except Exception as e:
            st.warning(f"OpenAI error: {e}. Using fallback model.")
            return fallback_titles()
    else:
        return fallback_titles()

def generate_slide_content(title, openai_api_key=None):
    """
    Generate PowerPoint slide content using OpenAI or HuggingFace
    """
    try:
        if openai_api_key:
            import openai
            client = openai.OpenAI(api_key=openai_api_key)
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a PowerPoint slide content generator."},
                    {"role": "user", "content": f"Write a professional 4-5 sentence paragraph for PowerPoint about: {title}"}
                ],
                temperature=0.7,
                max_tokens=200
            )
            return response.choices[0].message.content.strip()
        else:
            prompt = f"Generate a professional PowerPoint paragraph about {title} (4-5 sentences):"
            result = generator(
                prompt,
                max_new_tokens=120,
                num_return_sequences=1,
                truncation=True,
                pad_token_id=50256
            )[0]['generated_text']
            content = result.split(":")[-1].strip()
            return content
    except Exception as e:
        st.error(f"Error generating content for '{title}': {e}")
        return "(Content generation failed)"

def create_ppt(topic, titles, contents, bg_color="#FFFFFF", title_color="#000000", content_color="#000000", title_size=30, content_size=16, border_color=None, border_width=0):
    from pptx.dml.color import RGBColor
    prs = Presentation()
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = topic
    # Set title slide background color
    fill = title_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(int(bg_color[1:3], 16), int(bg_color[3:5], 16), int(bg_color[5:7], 16))
    # Set title color and size
    for p in title_slide.shapes.title.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(title_size)
            run.font.color.rgb = RGBColor(int(title_color[1:3], 16), int(title_color[3:5], 16), int(title_color[5:7], 16))
    # Content Slides
    content_slide_layout = prs.slide_layouts[1]  # Title and Content
    for title, content in zip(titles, contents):
        slide = prs.slides.add_slide(content_slide_layout)
        # Set background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(int(bg_color[1:3], 16), int(bg_color[3:5], 16), int(bg_color[5:7], 16))
        slide.shapes.title.text = title
        # Set title color and size
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(title_size)
                run.font.color.rgb = RGBColor(int(title_color[1:3], 16), int(title_color[3:5], 16), int(title_color[5:7], 16))
        # Add content to the body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Body placeholder
                shape.text = content
                # Set font size and color for all paragraphs
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(content_size)
                        run.font.color.rgb = RGBColor(int(content_color[1:3], 16), int(content_color[3:5], 16), int(content_color[5:7], 16))
                # Add border if requested
                if border_color and border_width > 0:
                    sp = shape
                    ln = sp.line
                    ln.color.rgb = RGBColor(int(border_color[1:3], 16), int(border_color[3:5], 16), int(border_color[5:7], 16))
                    ln.width = Pt(border_width)
        # Add default image quickly (local file)
        image_path = os.path.join(os.path.dirname(__file__), "default.jpg")
        if os.path.exists(image_path):
            # Insert image at a fixed position and size (customize as needed)
            left = Pt(400)
            top = Pt(100)
            width = Pt(200)
            slide.shapes.add_picture(image_path, left, top, width=width)
    os.makedirs("generated_ppt", exist_ok=True)
    safe_topic = "".join(c for c in topic if c.isalnum() or c in (" ", "_", "-"))
    path = f"generated_ppt/{safe_topic}_presentation.pptx"
    prs.save(path)
    return path

def get_download_link(path):
    with open(path, "rb") as f:
        data = f.read()
    b64 = base64.b64encode(data).decode()
    filename = os.path.basename(path)
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">📥 Download the PowerPoint</a>'

def main():
    st.title("Free Text-to-PPT Generator (CHATGPT)")
    st.write("Generate a professional PowerPoint presentation from a topic using AI.")
    topic = st.text_input("Enter your presentation topic:")
    openai_api_key = st.text_input("OpenAI API Key (optional, for better quality):", type="password")
    st.markdown("---")
    st.subheader("Customize Appearance")
    bg_color = st.color_picker("Slide Background Color", "#FFFFFF")
    title_color = st.color_picker("Title Text Color", "#000000")
    content_color = st.color_picker("Content Text Color", "#000000")
    title_size = st.slider("Title Font Size", 20, 60, 30)
    content_size = st.slider("Content Font Size", 10, 40, 16)
    border_color = st.color_picker("Border Color (optional)", "#000000")
    border_width = st.slider("Border Thickness (pt, 0 for none)", 0, 10, 0)
    if st.button("Generate Presentation"):
        if not topic:
            st.warning("Please enter a topic.")
            return
        st.info("Generating slides with AI...")
        titles = generate_slide_titles(topic, openai_api_key)
        if not titles or len(titles) < 1:
            st.error("Failed to generate slide titles. Please try again or check your API key.")
            return
        contents = [generate_slide_content(title, openai_api_key) for title in titles]
        path = create_ppt(topic, titles, contents, bg_color, title_color, content_color, title_size, content_size, border_color, border_width)
        st.success("✅ Presentation created!")
        st.markdown(get_download_link(path), unsafe_allow_html=True)

if __name__ == "__main__":
    main()
